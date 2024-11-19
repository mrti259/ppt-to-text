#!/usr/bin/env python
# coding: utf-8

# In[28]:


import os
from pathlib import Path

path_in = Path("in")
path_out = Path("out")

if not path_out.exists():
    os.mkdir(path_out)

ppt_in = [file for file in path_in.glob("**/*.pptx")]
list(ppt_in)


# In[42]:


from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract(path_ppt):
    prs = Presentation(path_ppt)
    lines = []
    
    dir_out = path_out / path_ppt.relative_to(path_in).with_suffix("")
    if dir_out.exists():
        return
        
    os.makedirs(dir_out)
    
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                filename = f"{i}-{j}_{shape.image.filename}"
                img_out = dir_out / filename
                if not img_out.exists():
                    with open(img_out, "wb") as file:
                        file.write(shape.image.blob)
                lines.append(f"![img]({filename})")
                continue
            
            if not shape.has_text_frame:
                continue
                
            for paragraph in shape.text_frame.paragraphs:
                line = ""
                for run in paragraph.runs:
                    text = run.text.strip()
                    if text == "": continue;
                    font = run.font
                    if font.bold: text = f"**{text}**";
                    if font.italic: text = f"*{text}*";
                    if font.underline: text = f"<u>{text}</u>";
                    if line != "": line += " ";
                    line += text
                if line == "": continue;
                lines.append(line)
                    
    with open(dir_out / "README.md", "w") as file:
        file.write("\n\n".join(lines))
        print(dir_out)


# In[43]:


for path_ppt in ppt_in:
    extract(path_ppt)

